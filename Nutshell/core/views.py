# myapp/views.py
import json
from rest_framework.decorators import api_view
from rest_framework.response import Response
from rest_framework import status
import nltk # type: ignore
from sentence_transformers import SentenceTransformer, util # type: ignore
import numpy as np # type: ignore
import logging
from django.views.decorators.csrf import csrf_exempt
from scipy.special import softmax # type: ignore
from scipy.sparse.csgraph import connected_components # type: ignore
from googletrans import Translator # type: ignore
################################################################
import os
from rest_framework.views import APIView
from rest_framework.parsers import MultiPartParser, FormParser
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
from pptx.util import Inches, Pt # type: ignore
from django.http import HttpResponse



#+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
#Summarize start

def degree_centrality_scores(similarity_matrix, threshold=None, increase_power=True):
    if not (threshold is None or isinstance(threshold, float) and 0 <= threshold < 1):
        raise ValueError("'threshold' should be a floating-point number from the interval [0, 1) or None")
    
    if threshold is None:
        markov_matrix = create_markov_matrix(similarity_matrix)
    else:
        markov_matrix = create_markov_matrix_discrete(similarity_matrix, threshold)
    
    scores = stationary_distribution(markov_matrix, increase_power=increase_power, normalized=False)
    return scores

def _power_method(transition_matrix, increase_power=True, max_iter=10000):
    eigenvector = np.ones(len(transition_matrix))
    if len(eigenvector) == 1:
        return eigenvector
    
    transition = transition_matrix.transpose()
    for _ in range(max_iter):
        eigenvector_next = np.dot(transition, eigenvector)
        if np.allclose(eigenvector_next, eigenvector):
            return eigenvector_next
        eigenvector = eigenvector_next
        if increase_power:
            transition = np.dot(transition, transition)
    return eigenvector_next

def connected_nodes(matrix):
    _, labels = connected_components(matrix)
    groups = []
    for tag in np.unique(labels):
        group = np.where(labels == tag)[0]
        groups.append(group)
    return groups

def create_markov_matrix(weights_matrix):
    n_1, n_2 = weights_matrix.shape
    if n_1 != n_2:
        raise ValueError("'weights_matrix' should be square")
    row_sum = weights_matrix.sum(axis=1, keepdims=True)
    if np.min(weights_matrix) <= 0:
        return softmax(weights_matrix, axis=1)
    return weights_matrix / row_sum

def create_markov_matrix_discrete(weights_matrix, threshold):
    discrete_weights_matrix = np.zeros(weights_matrix.shape)
    ixs = np.where(weights_matrix >= threshold)
    discrete_weights_matrix[ixs] = 1
    return create_markov_matrix(discrete_weights_matrix)

def stationary_distribution(transition_matrix, increase_power=True, normalized=True):
    n_1, n_2 = transition_matrix.shape
    if n_1 != n_2:
        raise ValueError("'transition_matrix' should be square")
    distribution = np.zeros(n_1)
    grouped_indices = connected_nodes(transition_matrix)
    for group in grouped_indices:
        t_matrix = transition_matrix[np.ix_(group, group)]
        eigenvector = _power_method(t_matrix, increase_power=increase_power)
        distribution[group] = eigenvector
    if normalized:
        distribution /= n_1
    return distribution

# Main summarization function
@api_view(["POST"])
@csrf_exempt
def Summarization(request):
    nltk.download('punkt')
    
    data = json.loads(request.body)
    document = data.get('document', '')

        # Check if document is provided
    if not document:
        return Response({"error": "No document provided."}, status=400)

    # Split the document into sentences
    sentences = nltk.sent_tokenize(document)
    print("Num sentences:", len(sentences))

    # Load the sentence transformer model
    model = SentenceTransformer("all-MiniLM-L6-v2")

    # Parse the input document from the request data
    document = request.data.get('document', '')
    
    if not document:
        return Response({"error": "No document provided"}, status=status.HTTP_400_BAD_REQUEST)

    # Split the document into sentences
    sentences = nltk.sent_tokenize(document)
    
    # Compute the sentence embeddings
    embeddings = model.encode(sentences, convert_to_tensor=True)

    # Compute the pair-wise cosine similarities
    cos_scores = util.cos_sim(embeddings, embeddings).numpy()

    # Compute the centrality for each sentence
    centrality_scores = degree_centrality_scores(cos_scores, threshold=None)

    # Sort the sentences by centrality scores
    most_central_sentence_indices = np.argsort(-centrality_scores)

    # Extract the top 5 most central sentences
    summary_sentences = [sentences[idx] for idx in most_central_sentence_indices[:5]]

    # Join the summary sentences into a single summary string
    summary = ' '.join(summary_sentences)

    return Response({"summary": summary}, status=status.HTTP_200_OK)

#Summarize End
#+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++



#+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
#Translate start

@api_view(['POST'])
@csrf_exempt
def translate_english_to_arabic(request):
    try:
        data = request.data
        text = data.get('text', '')
        if not text:
            return Response({'error': 'Text field is required.'}, status=status.HTTP_400_BAD_REQUEST)
        

        translator = Translator()
        translated = translator.translate(text, src='en', dest='ar')
        return Response({'translated_text': translated.text}, status=status.HTTP_200_OK)
    except Exception as e:
        return Response({'error': str(e)}, status=status.HTTP_400_BAD_REQUEST)
    
# Translate End
#+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++


# Book Start
#+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++

model = SentenceTransformer("all-MiniLM-L6-v2")

# Define the path to the Downloads directory
DOWNLOADS_PATH = os.path.join(os.path.expanduser('~'), 'Downloads')

class DocumentUploadView(APIView):
    parser_classes = (MultiPartParser, FormParser)

    def post(self, request, *args, **kwargs):
        if 'file' not in request.FILES:
            return Response({"error": "No file provided"}, status=status.HTTP_400_BAD_REQUEST)
        
        file = request.FILES['file']
        # Save the file to the Downloads directory
        file_name = 'uploaded_file.docx'
        temp_file_path = os.path.join(DOWNLOADS_PATH, file_name)
        with open(temp_file_path, 'wb+') as temp_file:
            for chunk in file.chunks():
                temp_file.write(chunk)

        try:
            # Process the file and generate the PowerPoint
            sections = self.split_chapter_into_sections(temp_file_path)
            ppt_file_name = 'generated_presentation.pptx'
            ppt_file_path = os.path.join(DOWNLOADS_PATH, ppt_file_name)
            self.save_summaries_to_pptx(sections, ppt_file_path)
            response = self.create_ppt_response(ppt_file_path)
            return response
        except Exception as e:
            return Response({"error": str(e)}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    def split_chapter_into_sections(self, docx_file_path):
        document = Document(docx_file_path)
        sections = []
        current_section = []
        current_header = None

        for para in document.paragraphs:
            stripped_line = para.text.strip()
            if stripped_line.isupper() and stripped_line:
                if current_section:
                    sections.append((current_header, '\n'.join(current_section)))
                    current_section = []
                current_header = stripped_line
            else:
                if current_header:
                    current_section.append(para.text)

        if current_section:
            sections.append((current_header, '\n'.join(current_section)))

        return sections

    def summarize_text(self, text):
        sentences = nltk.sent_tokenize(text)
        if len(sentences) == 0:
            return ["No content to summarize."]
        if len(sentences) == 1:
            return [sentences[0]]

        embeddings = model.encode(sentences, convert_to_tensor=True)
        cos_scores = util.cos_sim(embeddings, embeddings).numpy()
        centrality_scores = degree_centrality_scores(cos_scores, threshold=None)
        most_central_sentence_indices = np.argsort(-centrality_scores)
        summary = [sentences[idx].strip() for idx in most_central_sentence_indices[:5]]

        return summary

    def save_summaries_to_pptx(self, sections, output_pptx_file_path):
        prs = Presentation()
        for i, (header, section) in enumerate(sections):
            summary = self.summarize_text(section)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = header if header else f"Section {i+1} Summary"

            content = '\n'.join([f"- {bullet_point}" for bullet_point in summary])
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for bullet_point in summary:
                p = text_frame.add_paragraph()
                p.text = bullet_point
                p.space_after = Inches(0.1)
                p.level = 0
                p.font.name = 'Arial'
                p.font.size = Pt(12)

        prs.save(output_pptx_file_path)

    def create_ppt_response(self, ppt_file_path):
        with open(ppt_file_path, 'rb') as ppt_file:
            response = HttpResponse(ppt_file.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename="{os.path.basename(ppt_file_path)}"'
            return response

# Book End
#+++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++++
